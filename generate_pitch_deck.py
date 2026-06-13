"""
Chapter 2 Drip — PREMIUM BRIGHT D2C Pitch Deck
Palette: White/Off-White BG + Slate-900 Type + Electric accents (Indigo, Lime, Rose)
Based on real Meta Ads Library audit (June 2026) and ₹90L/month spend baseline
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── PREMIUM WARM STREETWEAR PALETTE (GOOGLE SLIDES THEME) ─────────────
C_BG         = RGBColor(245, 242, 235)   # Warm Cream / Off-White background
C_PANEL      = RGBColor(235, 230, 220)   # Slightly darker cream
C_CARD       = RGBColor(255, 255, 255)   # White card background
C_CARD2      = RGBColor(240, 236, 228)   # Alternating card background

# Forest Green / Navy / Charcoal / Terracotta / Gold Palette (Slide 14)
C_FOREST     = RGBColor(34,  76,  56)    # Forest Green
C_NAVY       = RGBColor(27,  54,  93)    # Deep Navy Blue
C_CHARCOAL   = RGBColor(15,  23,  42)    # Dark Slate / Black
C_TERRACOTTA = RGBColor(226, 123,  88)    # Rust / Terracotta Orange
C_GOLD       = RGBColor(235, 176,  85)    # Warm Gold Accent

C_WHITE      = RGBColor(255, 255, 255)   # True white
C_TEXT       = RGBColor(15,  23,  42)    # Slate 900 (Main dark text)
C_WD         = RGBColor(71,  85, 105)    # Slate 600 (Secondary text)
C_GREY       = RGBColor(100, 116, 139)   # Slate 500
C_GREY2      = RGBColor(148, 163, 184)   # Slate 400
C_BORDER     = RGBColor(218, 212, 198)   # Warm slate border

FONT_HDR     = "Trebuchet MS"
FONT_BODY    = "Arial"
W = Inches(13.333)
H = Inches(7.5)

# ── PRIMITIVES ───────────────────────────────────────────────────────
def bg(s, c=None):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c or C_BG

def box(s, l, t, w, h, fill, line=None, lw=0.6):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def rbox(s, l, t, w, h, fill, line=None, lw=0.8):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    sh.adjustments[0] = 0.04
    return sh

def ell(s, l, t, w, h, fill, line=None, lw=1.2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def bar(s, l, t, w, h=Inches(0.035), c=None):
    box(s, l, t, w, h, c or C_GOLD)

def arr(s, l, t, w=Inches(0.4), c=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.22))
    sh.fill.solid(); sh.fill.fore_color.rgb = c or C_GOLD
    sh.line.fill.background()

def txt(s, text, l, t, w, h, size=10, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True, ml=0, mt=0, font=None):
    color = color or C_TEXT
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = Inches(ml); tf.margin_top = Inches(mt)
    tf.margin_right = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font or FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color

def footer(s, n, total=8):
    box(s, 0, Inches(7.22), W, Inches(0.02), C_BORDER)
    txt(s, "CHAPTER 2 DRIP  ·  GROWTH & ROAS STRATEGY  ·  CONFIDENTIAL",
        Inches(0.4), Inches(7.26), W - Inches(1.2), Inches(0.2), size=7.5, color=C_GREY, font=FONT_BODY)
    txt(s, f"{n} / {total}", Inches(12.4), Inches(7.26), Inches(0.5), Inches(0.2),
        size=8, bold=True, color=C_FOREST, align=PP_ALIGN.RIGHT, font=FONT_BODY)

def hdr(s, title, sub=None, accent=None):
    # Solid Forest Green Header banner (just like slide 14)
    box(s, 0, 0, W, Inches(1.16), C_FOREST)
    # Thin Accent Gold bar at the bottom
    box(s, 0, Inches(1.16), W, Inches(0.04), C_GOLD)
    # Title in white
    txt(s, title.upper(), Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.6),
        size=20, bold=True, color=C_WHITE, ml=0, font=FONT_HDR)
    if sub:
        # Subtitle in warm gold/cream
        txt(s, sub, Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.38),
            size=9.5, italic=True, color=C_GOLD, ml=0, font=FONT_BODY)

def mcard(s, l, t, w, h, label, val, note="", vc=None, fc=None, ac=None, vsz=24):
    vc = vc or C_GOLD; fc = fc or C_CARD; ac = ac or C_GOLD
    rbox(s, l, t, w, h, fc, ac, 0.8)
    bar(s, l, t, w, Inches(0.045), ac)
    txt(s, label, l, t + Inches(0.09), w, Inches(0.28),
        size=7.5, bold=True, color=C_GREY, align=PP_ALIGN.CENTER, ml=0, font=FONT_BODY)
    txt(s, val, l, t + Inches(0.32), w, Inches(0.56),
        size=vsz, bold=True, color=vc, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
    if note:
        txt(s, note, l, t + h - Inches(0.36), w, Inches(0.3),
            size=8, italic=True, color=C_GREY, align=PP_ALIGN.CENTER, ml=0, font=FONT_BODY)

def badge(s, l, t, num, col):
    ell(s, l, t, Inches(0.38), Inches(0.38), col)
    txt(s, num, l, t + Inches(0.07), Inches(0.38), Inches(0.28),
        size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)

# ── SLIDE 1: COVER ───────────────────────────────────────────────────
def s01(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    box(s, 0, 0, Inches(0.2), H, C_FOREST)
    box(s, Inches(0.2), 0, Inches(0.05), H, C_GOLD)
    
    # Left text content
    txt(s, "CHAPTER 2", Inches(0.6), Inches(1.3), Inches(7.5), Inches(1.05),
        size=55, bold=True, color=C_FOREST, font=FONT_HDR)
    txt(s, "DRIP.", Inches(0.6), Inches(2.25), Inches(7.5), Inches(1.05),
        size=55, bold=True, color=C_TERRACOTTA, font=FONT_HDR)
    bar(s, Inches(0.6), Inches(3.45), Inches(5.0), Inches(0.05), C_FOREST)
    
    txt(s, "GROWTH STRATEGY  ·  CEO BRIEFING",
        Inches(0.6), Inches(3.62), Inches(7.5), Inches(0.42),
        size=13, bold=True, color=C_FOREST, font=FONT_HDR)
    txt(s, "Audit-backed roadmap to turn ₹90L/month in ad spend\ninto a profitable, compounding growth engine.",
        Inches(0.6), Inches(4.14), Inches(7.5), Inches(0.8),
        size=12, italic=True, color=C_WD, font=FONT_BODY)
    
    # Right-side solid Navy panel for stats
    box(s, Inches(8.5), Inches(0.4), Inches(4.4), Inches(6.6), C_NAVY)
    box(s, Inches(8.5), Inches(0.4), Inches(4.4), Inches(0.06), C_GOLD)
    
    stats = [
        ("CURRENT AD SPEND",  "₹90L/mo",   C_WHITE),
        ("CURRENT AD ROAS",   "1.8×",       C_TERRACOTTA),
        ("TARGET AD ROAS",    "3.2×",       C_GOLD),
        ("TARGET REVENUE",    "₹4.32 Cr",  C_WHITE),
    ]
    sy = 0.8
    for lbl, val, col in stats:
        txt(s, lbl, Inches(8.9), Inches(sy), Inches(3.6), Inches(0.28),
            size=8.5, bold=True, color=C_GOLD, ml=0, font=FONT_HDR)
        txt(s, val, Inches(8.9), Inches(sy + 0.26), Inches(3.6), Inches(0.46),
            size=22, bold=True, color=col, ml=0, font=FONT_HDR)
        box(s, Inches(8.9), Inches(sy + 0.78), Inches(3.6), Inches(0.015), C_TERRACOTTA)
        sy += 1.15
        
    box(s, Inches(8.9), Inches(6.0), Inches(3.6), Inches(0.7), C_CHARCOAL)
    txt(s, "Strategic Partner: Marin Istanovic\nMedia Model: Creative-Led Broad Targeting",
        Inches(8.9), Inches(6.05), Inches(3.6), Inches(0.6),
        size=8.5, italic=True, color=C_WHITE, ml=0.1, mt=0.04, font=FONT_BODY)
        
    footer(s, 1)

# ── SLIDE 2: ADS AUDIT ──────────────────────────────────────────────
def s02(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "Meta Ads Library Audit — Chapter 2 Drip",
        "What the brand is actually running right now — and what's broken in the performance engine.", C_TERRACOTTA)
    footer(s, 2)
    
    # Solid dark-colored cards matching the Google Slide theme columns
    obs = [
        (C_FOREST, "WHAT THEY'RE RUNNING",
         "Bundle pack deals (Trio Pack ₹5,780 · Cargo combos ₹2,999)\n"
         "Earth-tone colorways — heavy focus on Olive, Rust, Brown\n"
         "7 variations of the SAME copy running simultaneously\n"
         "Price anchoring ('₹3,990 Combo') used as the main hook"),
        (C_NAVY, "HOW THEY'RE SAYING IT",
         "'Earthy fix with the Olive Raasta & Brown Munde Pack'\n"
         "'Utility-inspired cargo with washed olive finish'\n"
         "'Built for comfort and style — works for any occasion'\n"
         "Product specifications as copy. No hook, no emotion."),
        (C_CHARCOAL, "WHAT FORMAT / CREATIVE",
         "Static flat-lays & model images — zero video detected\n"
         "Multiple ad variations testing the exact same headline\n"
         "Direct price CTA: 'Shop Now for ₹3,990'\n"
         "No UGC, no creator content, no streetwear lifestyle"),
    ]
    ox = Inches(0.4)
    for col, title, body in obs:
        # Solid card background
        rbox(s, ox, Inches(1.3), Inches(4.05), Inches(2.3), col)
        bar(s, ox, Inches(1.3), Inches(4.05), Inches(0.05), C_GOLD)
        txt(s, title, ox + Inches(0.15), Inches(1.4), Inches(3.75), Inches(0.3),
            size=9, bold=True, color=C_GOLD, ml=0, font=FONT_HDR)
        txt(s, body, ox + Inches(0.15), Inches(1.78), Inches(3.75), Inches(1.7),
            size=9.5, color=C_WHITE, wrap=True, ml=0, font=FONT_BODY)
        ox += Inches(4.25)
        
    box(s, Inches(0.4), Inches(3.72), W - Inches(0.8), Inches(0.02), C_BORDER)
    txt(s, "6 CRITICAL PERFORMANCE GAPS IN CURRENT AD SYSTEM", Inches(0.4), Inches(3.82),
        Inches(8.0), Inches(0.3), size=9.5, bold=True, color=C_FOREST, ml=0, font=FONT_HDR)
        
    gaps = [
        (C_FOREST, "NO HOOK STRATEGY", "Ad copy opens with product name/color, not an engaging hook."),
        (C_TERRACOTTA, "7 VERSIONS, 1 IDEA", "Testing volume rather than core creative/messaging diversity."),
        (C_NAVY, "ZERO UGC / fit-checks", "No creator fit-checks, styling videos, or product hauls."),
        (C_CHARCOAL, "NO SCARCITY MECHANICS", "Ads run indefinitely with no drop FOMO or urgency hooks."),
        (C_TERRACOTTA, "NO VIDEO CREATIVE", "Zero video ads or motion graphics running in active library."),
        (C_FOREST, "NO TRUST SIGNALS", "No customer reviews, badges, or social proof on creatives."),
    ]
    gx = Inches(0.4); gy = Inches(4.18)
    for i, (col, title, body) in enumerate(gaps):
        lx = gx + (i % 3) * Inches(4.25)
        ty = gy + (i // 3) * Inches(1.42)
        rbox(s, lx, ty, Inches(4.05), Inches(1.32), C_CARD, C_BORDER, 0.7)
        bar(s, lx, ty, Inches(4.05), Inches(0.045), col)
        txt(s, title, lx + Inches(0.15), ty + Inches(0.1), Inches(3.75), Inches(0.26),
            size=8.5, bold=True, color=col, ml=0, font=FONT_HDR)
        txt(s, body, lx + Inches(0.15), ty + Inches(0.42), Inches(3.75), Inches(0.82),
            size=9, color=C_TEXT, wrap=True, ml=0, font=FONT_BODY)

# ── SLIDE 3: CREATIVE ENGINE ──────────────────────────────────────────
def s03(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "The Fix — Creative-Led Broad Targeting",
        "One campaign. No interest or location stacking. The creative IS the targeting signal.", C_NAVY)
    footer(s, 3)
    
    # Sub-header quote
    rbox(s, Inches(0.4), Inches(1.26), W - Inches(0.8), Inches(0.56), C_CARD2, C_NAVY, 0.8)
    bar(s, Inches(0.4), Inches(1.26), Inches(0.07), Inches(0.56), C_NAVY)
    txt(s, '\u201cStop targeting interests. Start targeting with creatives. Broad + Advantage+ Shopping Campaign.\u201d — Marin Istanovic',
        Inches(0.65), Inches(1.37), W - Inches(1.2), Inches(0.38),
        size=10, italic=True, color=C_NAVY, ml=0, font=FONT_HDR)
        
    steps = [
        ("1  BRIEF",   "Pick ONE angle.",     C_FOREST),
        ("2  PRODUCE", "10–15 UGC clips.",      C_NAVY),
        ("3  LAUNCH",  "1 ASC+ campaign.", C_CHARCOAL),
        ("4  JUDGE",   "72-hour rule.",          C_TERRACOTTA),
        ("5  SCALE",   "Winners get 5\u00d7.",    C_FOREST),
    ]
    nw, nh = Inches(2.2), Inches(1.52)
    for i, (title, sub, col) in enumerate(steps):
        lx = Inches(0.4) + i * (nw + Inches(0.26))
        rbox(s, lx, Inches(2.02), nw, nh, C_CARD, col, 1.2)
        bar(s, lx, Inches(2.02), nw, Inches(0.055), col)
        txt(s, title, lx, Inches(2.14), nw, Inches(0.4),
            size=12, bold=True, color=col, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        txt(s, sub, lx, Inches(2.6), nw, Inches(0.8),
            size=9.5, color=C_TEXT, align=PP_ALIGN.CENTER, wrap=True, ml=0, font=FONT_BODY)
        if i < 4: arr(s, lx + nw + Inches(0.05), Inches(2.68), Inches(0.16), C_GOLD)
        
    # Bottom elements
    rbox(s, Inches(0.4), Inches(3.72), Inches(5.9), Inches(3.2), C_CARD, C_TERRACOTTA, 0.8)
    bar(s, Inches(0.4), Inches(3.72), Inches(5.9), Inches(0.055), C_TERRACOTTA)
    txt(s, "THE 72-HOUR PERFORMANCE RULE", Inches(0.6), Inches(3.86), Inches(5.5), Inches(0.3),
        size=9.5, bold=True, color=C_TERRACOTTA, ml=0, font=FONT_HDR)
        
    kill = [
        ("Hook Rate < 25%",           "Kill immediately",        C_TERRACOTTA),
        ("Spend > \u20b93K, ROAS < 2.0",  "Kill",                    C_TERRACOTTA),
        ("CPM \u2191 + CTR \u2193",            "Rewrite the hook",        C_TERRACOTTA),
        ("ROAS 2.0\u20133.0",              "Keep \u2014 5 more days", C_GOLD),
        ("ROAS > 3.2\u00d7",              "Scale budget 5\u00d7",    C_FOREST),
    ]
    ky = Inches(4.28)
    for cond, action, col in kill:
        box(s, Inches(0.6), ky, Inches(5.5), Inches(0.4), C_CARD2, C_BORDER, 0.3)
        txt(s, cond, Inches(0.72), ky + Inches(0.08), Inches(3.1), Inches(0.26),
            size=9, color=C_TEXT, ml=0, font=FONT_BODY)
        txt(s, action, Inches(3.85), ky + Inches(0.08), Inches(2.1), Inches(0.26),
            size=9, bold=True, color=col, ml=0, font=FONT_HDR)
        ky += Inches(0.48)
        
    rbox(s, Inches(6.65), Inches(3.72), Inches(6.28), Inches(3.2), C_CARD, C_FOREST, 0.8)
    bar(s, Inches(6.65), Inches(3.72), Inches(6.28), Inches(0.055), C_FOREST)
    txt(s, "7 NARRATIVE ANGLES TO TEST", Inches(6.85), Inches(3.86), Inches(5.8), Inches(0.3),
        size=9.5, bold=True, color=C_FOREST, ml=0, font=FONT_HDR)
        
    angles = [
        (C_FOREST,     "#1  Scarcity",   "'500 units · Drop 02 · Never Restocked.'"),
        (C_NAVY,       "#2  Spec-led",   "'240 GSM heavy French Terry · Puff Print.'"),
        (C_TERRACOTTA, "#3  Combo Pack", "'Trio Pack — get any 3 shirts & save ₹1,790.'"),
        (C_GOLD,       "#4  Comfort",    "'Worn once. Immediately ordered again.'"),
        (C_CHARCOAL,   "#5  Identity",   "'Not everyone gets the backprint reference.'"),
        (C_NAVY,       "#6  Haul UGC",    "Creator styling & fit-check review video."),
        (C_TERRACOTTA, "#7  Hinglish",   "'Drip level max.' Mumbai streetwear hooks."),
    ]
    ay = Inches(4.24)
    for col, label, copy in angles:
        ell(s, Inches(6.85), ay + Inches(0.05), Inches(0.2), Inches(0.2), col)
        txt(s, label, Inches(7.15), ay, Inches(1.4), Inches(0.3),
            size=8.5, bold=True, color=col, ml=0, font=FONT_HDR)
        txt(s, copy, Inches(8.65), ay, Inches(4.1), Inches(0.3),
            size=8.5, italic=True, color=C_TEXT, ml=0, font=FONT_BODY)
        ay += Inches(0.39)

# ── SLIDE 4: RETENTION ───────────────────────────────────────────────
def s04(prs, blank):
    import math
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "Retention \u2014 The Revenue That Compounds",
        "Acquiring a customer costs \u20b91,222. Keeping one costs only \u20b9240.", C_FOREST)
    footer(s, 4)
    cx, cy = 4.6, 4.4
    
    # Outer circle guide
    ell(s, Inches(cx-1.55), Inches(cy-1.55), Inches(3.1), Inches(3.1), C_CARD2, C_FOREST, 0.7)
    # Inner green circle
    ell(s, Inches(cx-0.95), Inches(cy-0.95), Inches(1.9), Inches(1.9), C_FOREST, C_GOLD, 1.8)
    txt(s, "GROWTH\nFLYWHEEL", Inches(cx-0.95), Inches(cy-0.56),
        Inches(1.9), Inches(0.9), size=11, bold=True, color=C_GOLD,
        align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        
    orbit = [
        ("ACQUIRE",    "1st buy",         C_FOREST),
        ("NURTURE",    "Post-buy flow",   C_NAVY),
        ("RETAIN",     "2nd buy",         C_CHARCOAL),
        ("VIP ACCESS", "Early drop",      C_TERRACOTTA),
        ("REFERRAL",   "WOM Gift",        C_FOREST),
        ("WIN-BACK",   "Day 60 nudge",    C_TERRACOTTA),
    ]
    r_orb = 2.65
    for i, (lbl, sub, col) in enumerate(orbit):
        angle = math.radians(90 + i * 60)
        ox = cx + r_orb * math.cos(angle)
        oy = cy - r_orb * math.sin(angle)
        ow, oh = 1.35, 0.88
        rbox(s, Inches(ox-ow/2), Inches(oy-oh/2), Inches(ow), Inches(oh), C_CARD, col, 1.2)
        bar(s, Inches(ox-ow/2), Inches(oy-oh/2), Inches(ow), Inches(0.04), col)
        txt(s, lbl, Inches(ox-ow/2), Inches(oy-oh/2+0.06), Inches(ow), Inches(0.3), size=8.5, bold=True, color=col, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        txt(s, sub, Inches(ox-ow/2), Inches(oy-oh/2+0.38), Inches(ow), Inches(0.3), size=7.5, color=C_TEXT, align=PP_ALIGN.CENTER, ml=0, font=FONT_BODY)
        
    # Calendar panel on the right
    rbox(s, Inches(9.8), Inches(1.26), Inches(3.2), Inches(5.82), C_CARD, C_FOREST, 0.8)
    bar(s, Inches(9.8), Inches(1.26), Inches(3.2), Inches(0.055), C_FOREST)
    txt(s, "RETENTION CALENDAR", Inches(10.0), Inches(1.4), Inches(2.8), Inches(0.3), size=9.5, bold=True, color=C_FOREST, ml=0, font=FONT_HDR)
    
    cal = [
        ("Day 0",  "Welcome SMS flow",      C_FOREST),
        ("Day 3",  "Streetwear style guide", C_NAVY),
        ("Day 7",  "Review & UGC incentive", C_CHARCOAL),
        ("Day 14", "Coord cross-sell pack",  C_TERRACOTTA),
        ("Day 30", "Next collection VIP drop", C_GOLD),
        ("Day 60", "Lapsed buyer win-back",  C_TERRACOTTA),
    ]
    ty2 = Inches(1.86)
    for day, action, col in cal:
        ell(s, Inches(10.0), ty2 + Inches(0.04), Inches(0.22), Inches(0.22), col)
        txt(s, day, Inches(10.32), ty2, Inches(0.7), Inches(0.3), size=8.5, bold=True, color=col, ml=0, font=FONT_HDR)
        txt(s, action, Inches(11.1), ty2, Inches(1.8), Inches(0.32), size=8.5, color=C_TEXT, wrap=True, ml=0, font=FONT_BODY)
        ty2 += Inches(0.62)

# ── SLIDE 5: ECONOMICS ───────────────────────────────────────────────
def s05(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "The Business Case \u2014 Before vs. After",
        "Same \u20b990L media spend baseline. Smarter operations. Radically different profitability outcome.", C_GOLD)
    footer(s, 5)

    cols = [
        ("TODAY (PR-LED)", C_CHARCOAL, [
            ("Monthly Media Spend", "₹90L",      "/ month", C_WHITE),
            ("Ad Spend ROAS",     "1.8\u00d7",   "",        C_TERRACOTTA),
            ("Gross Revenue",     "₹1.62 Cr",  "/ month", C_WHITE),
            ("Average Order Value","₹2,200",    "",        C_WHITE),
            ("Contribution Margin","\u2013\u20b913.5L", "(–8.3% CM1)", C_TERRACOTTA),
        ]),
        ("90-DAY TARGET", C_FOREST, [
            ("Monthly Media Spend", "₹1.35 Cr",  "/ month", C_WHITE),
            ("Ad Spend ROAS",     "3.2\u00d7",   "",        C_GOLD),
            ("Gross Revenue",     "₹4.32 Cr",  "/ month", C_WHITE),
            ("Average Order Value","₹3,200",    "",        C_WHITE),
            ("Contribution Margin","+\u20b91.08 Cr", "(+25.0% CM1)", C_GOLD),
        ]),
    ]
    for ci, (label, col, rows) in enumerate(cols):
        lx = Inches(0.4 + ci * 5.65)
        cw = Inches(5.2)
        rbox(s, lx, Inches(1.28), cw, Inches(5.82), col)
        bar(s, lx, Inches(1.28), cw, Inches(0.06), C_GOLD)
        txt(s, label, lx, Inches(1.38), cw, Inches(0.38),
            size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        ry = Inches(1.82)
        for metric, val, unit, vc in rows:
            # Subtle card nested in solid card
            box(s, lx + Inches(0.18), ry, cw - Inches(0.36), Inches(0.92),
                C_CHARCOAL if col == C_CHARCOAL else C_FOREST, C_BORDER, 0.4)
            txt(s, metric, lx + Inches(0.28), ry + Inches(0.06),
                cw - Inches(0.5), Inches(0.26), size=8.5, color=C_GOLD, ml=0, font=FONT_HDR)
            txt(s, val, lx + Inches(0.28), ry + Inches(0.34),
                Inches(3.4), Inches(0.45), size=18, bold=True, color=vc, ml=0, font=FONT_HDR)
            if unit:
                txt(s, unit, lx + cw - Inches(1.6), ry + Inches(0.48),
                    Inches(1.3), Inches(0.26), size=8, italic=True, color=C_GOLD,
                    align=PP_ALIGN.RIGHT, ml=0, font=FONT_BODY)
            ry += Inches(0.96)

    arr(s, Inches(5.72), Inches(4.05), Inches(0.45), C_GOLD)
    txt(s, "90\nDAYS", Inches(5.7), Inches(3.72), Inches(0.48), Inches(0.5),
        size=7.5, bold=True, color=C_CHARCOAL, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)

    # Levers list on the right
    rbox(s, Inches(11.35), Inches(1.28), Inches(1.7), Inches(5.82), C_CARD, C_BORDER, 0.5)
    txt(s, "GROWTH\nLEVERS", Inches(11.35), Inches(1.44), Inches(1.7), Inches(0.55),
        size=9, bold=True, color=C_FOREST, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
    levers = [
        (C_FOREST,     "Creative\nVelocity"),
        (C_NAVY,       "AOV Packs\n& Bundles"),
        (C_CHARCOAL,   "Retention\nFlywheel"),
        (C_TERRACOTTA, "Drop\nScarcity"),
        (C_GOLD,       "Google\nBrand Search"),
    ]
    ly = Inches(2.1)
    for col, lbl in levers:
        ell(s, Inches(12.02), ly, Inches(0.32), Inches(0.32), col)
        txt(s, lbl, Inches(11.35), ly + Inches(0.38), Inches(1.7), Inches(0.44),
            size=7.5, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER, wrap=True, ml=0, font=FONT_BODY)
        ly += Inches(0.96)


# ── SLIDE 6: COMPETITOR CARDS ─────────────────────────────────────────
def s06(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "What Competitors Are Doing \u2014 Steal & Adapt",
        "Four brands doing parts of it right. Real ad intelligence. Clear actions for Chapter 2.", C_NAVY)
    footer(s, 6)

    comps = [
        {
            "name": "JAYWALKING", "tier": "Ultra-Premium · ₹3k–₹8k",
            "col": C_CHARCOAL, "border_col": C_GOLD, "tag": "CULT BRAND",
            "what": [
                "Zero discounts ever — high price anchors",
                "Silhouette-forward photography, minimal copy",
                "Designer & rapper seeding, zero paid influencers",
                "Very low creative velocity — scarcity is the hook",
            ],
            "steal": "Run product SILHOUETTE ads, not story ads. Let the design print and shape speak for itself.",
            "kpi": "Avg Selling Price", "kpi_v": "₹5,500+",
        },
        {
            "name": "BLUORNG", "tier": "Premium · ₹3.5k–₹7k",
            "col": C_TERRACOTTA, "border_col": C_GOLD, "tag": "DROP MODEL",
            "what": [
                "Never-restock model: drop countdowns in ads",
                "Fabric texture close-ups: GSM & puff print detail",
                "FOMO copy: 'This colorway will not return'",
                "Heavy Meta Advantage+ Catalog retargeting",
            ],
            "steal": "Add drop code to ads: 'Drop 02 · 500 units available · Never Restocked.'",
            "kpi": "Repeat Buy Rate", "kpi_v": "35%+",
        },
        {
            "name": "SNITCH", "tier": "Fast Fashion · ₹999–₹1.9k",
            "col": C_NAVY, "border_col": C_GOLD, "tag": "VOLUME SCALE",
            "what": [
                "Advantage+ Shopping Campaigns (ASC+) broad focus",
                "30+ new ad variations shipped weekly to the feed",
                "Fast-cut vertical reels & meme-based hooks",
                "Google PMax running alongside Meta broad scale",
            ],
            "steal": "Set up 1 ASC+ campaign. Partner with agency to ship 10 new UGC creatives weekly.",
            "kpi": "Weekly Creatives", "kpi_v": "30+",
        },
        {
            "name": "BONKERS CORNER", "tier": "Affordable · ₹799–₹1.5k",
            "col": C_FOREST, "border_col": C_GOLD, "tag": "UGC ENGINE",
            "what": [
                "College influencer hauls — low cost, high trust",
                "Relatable hooks: 'POV: your wallet after this haul'",
                "Heavy Snapchat story ads for younger cohorts",
                "Organic video virality boosted with paid ad spend",
            ],
            "steal": "Haul-format reels showing fits. Maintain premium vibes — no humor, just street aesthetics.",
            "kpi": "Monthly UGC Videos", "kpi_v": "50+",
        },
    ]

    cw, ch = Inches(3.05), Inches(5.82)
    gap = Inches(0.12)
    sx = Inches(0.4)

    for i, c in enumerate(comps):
        lx = sx + i * (cw + gap)
        col = c["col"]
        # Solid card panel
        rbox(s, lx, Inches(1.28), cw, ch, col)
        bar(s, lx, Inches(1.28), cw, Inches(0.06), c["border_col"])

        # Badge
        bw = Inches(1.1)
        box(s, lx + cw - bw - Inches(0.12), Inches(1.36), bw, Inches(0.26), c["border_col"])
        txt(s, c["tag"], lx + cw - bw - Inches(0.12), Inches(1.38), bw, Inches(0.24),
            size=7, bold=True, color=col, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)

        txt(s, c["name"], lx + Inches(0.15), Inches(1.36), cw - bw - Inches(0.3), Inches(0.34),
            size=11, bold=True, color=C_WHITE, ml=0, font=FONT_HDR)
        txt(s, c["tier"], lx + Inches(0.15), Inches(1.72), cw - Inches(0.3), Inches(0.24),
            size=7.5, italic=True, color=c["border_col"], ml=0, font=FONT_BODY)

        # AOV Metric nested card
        mcard(s, lx + Inches(0.15), Inches(2.02), cw - Inches(0.3), Inches(0.82),
              c["kpi"], c["kpi_v"], fc=col, vc=C_WHITE, ac=c["border_col"], vsz=15)

        fy = Inches(2.95)
        for fact in c["what"]:
            ell(s, lx + Inches(0.15), fy + Inches(0.06), Inches(0.18), Inches(0.18), c["border_col"])
            txt(s, fact, lx + Inches(0.42), fy, cw - Inches(0.55), Inches(0.36),
                size=8.5, color=C_WHITE, wrap=True, ml=0, font=FONT_BODY)
            fy += Inches(0.44)

        steal_y = Inches(1.28) + ch - Inches(1.1)
        box(s, lx, steal_y, cw, Inches(1.1), C_CHARCOAL if col != C_CHARCOAL else C_NAVY)
        bar(s, lx, steal_y, cw, Inches(0.035), c["border_col"])
        txt(s, "\u2192 STEAL & ADAPT FOR C2", lx + Inches(0.15), steal_y + Inches(0.08),
            cw - Inches(0.3), Inches(0.24), size=7.5, bold=True, color=c["border_col"], ml=0, font=FONT_HDR)
        txt(s, c["steal"], lx + Inches(0.15), steal_y + Inches(0.32),
            cw - Inches(0.3), Inches(0.7), size=8, color=C_WHITE, wrap=True, ml=0, font=FONT_BODY)


# ── SLIDE 7: ROADMAP ─────────────────────────────────────────────────
def s07(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    hdr(s, "90-Day Execution Roadmap",
        "Three focused sprints. One profitable outcome.", C_FOREST)
    footer(s, 7)

    phases = [
        ("MONTH 1", C_NAVY, "FIX THE ENGINE", [
            "Advantage+ broad focus - zero interest stack",
            "Onboard 20 streetwear creators (10 UGC/wk)",
            "Day-0 to Day-90 automated retention flows",
            "Google Brand Search campaign launched",
            "Shopify PDP checkout and bundle builder updates",
        ], "ASC+ live · Retention active · UGC pipeline running"),
        ("MONTH 2", C_TERRACOTTA, "SCALE WINNERS", [
            "5\u00d7 budget on ROAS > 3.2\u00d7 active creatives",
            "Kill all ad variations under 25% hook rate",
            "Meta catalog DPA dynamic carousels live",
            "VIP drop early-access email/SMS blasts",
            "PDP A/B tests: single item vs. bundle sets",
        ], "2-3 scaling creatives · AOV \u20b92,800+ · DPAs active"),
        ("MONTH 3", C_FOREST, "COMPOUND & RETAIN", [
            "Cross-sell SMS automation: Day-14 post-buy",
            "Referral engine: get \u20b9200, give \u20b9200",
            "Win-back automated loop for 60-day lapsed",
            "Test localized Hinglish reels hooks (Mumbai/Delhi)",
            "ROAS audit: blend search & paid at target \u2265 3.2\u00d7",
        ], "ROAS 3.2\u00d7 · CM1 +ve · Repeat 20%+ · LTV scaling"),
    ]

    tl_y = Inches(4.4)
    bar(s, Inches(0.45), tl_y, Inches(12.4), Inches(0.04), C_BORDER)

    pw, ph = Inches(4.02), Inches(2.95)
    for i, (month, col, theme, tasks, outcome) in enumerate(phases):
        lx = Inches(0.45) + i * (pw + Inches(0.18))
        rbox(s, lx, Inches(1.28), pw, ph, C_CARD, col, 0.9)
        bar(s, lx, Inches(1.28), pw, Inches(0.06), col)
        txt(s, month, lx, Inches(1.38), pw, Inches(0.3),
            size=8, bold=True, color=col, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        txt(s, theme, lx, Inches(1.72), pw, Inches(0.34),
            size=12.5, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        bar(s, lx + Inches(0.2), Inches(2.1), pw - Inches(0.4), Inches(0.025), C_BORDER)
        for j, task in enumerate(tasks):
            ell(s, lx + Inches(0.18), Inches(2.18) + j*Inches(0.38) + Inches(0.05),
                Inches(0.18), Inches(0.18), col)
            txt(s, task, lx + Inches(0.45), Inches(2.18) + j * Inches(0.38),
                pw - Inches(0.65), Inches(0.35), size=8.5, color=C_TEXT, ml=0, font=FONT_BODY)

        # Timeline dot
        dc = lx + pw / 2
        ell(s, dc - Inches(0.2), tl_y - Inches(0.2), Inches(0.4), Inches(0.4), col)
        txt(s, str(i+1), dc - Inches(0.2), tl_y - Inches(0.16),
            Inches(0.4), Inches(0.34), size=10, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)

        rbox(s, lx, tl_y + Inches(0.28), pw, Inches(0.65), C_CARD2, col, 0.5)
        txt(s, outcome, lx, tl_y + Inches(0.36), pw, Inches(0.52),
            size=7.5, italic=True, color=col, align=PP_ALIGN.CENTER, wrap=True, ml=0, font=FONT_BODY)


# ── SLIDE 8: CLOSE ────────────────────────────────────────────────────
def s08(prs, blank):
    s = prs.slides.add_slide(blank); bg(s, C_BG)
    box(s, 0, 0, Inches(0.2), H, C_FOREST)
    box(s, Inches(0.2), 0, Inches(0.05), H, C_GOLD)
    box(s, Inches(8.4), 0, W - Inches(8.4), H, C_PANEL)
    box(s, Inches(8.35), 0, Inches(0.05), H, C_TERRACOTTA)

    txt(s, "THE COMMITMENT",
        Inches(0.6), Inches(1.1), Inches(7.5), Inches(0.44),
        size=11, bold=True, color=C_GREY, font=FONT_HDR)
    txt(s, "Three deliverables for scale.",
        Inches(0.6), Inches(1.54), Inches(7.5), Inches(1.1),
        size=38, bold=True, color=C_TEXT, font=FONT_HDR)

    commits = [
        (C_FOREST,     "01", "A creative-led machine",      "No celebrity dependency. The creative matches and targets the buyer."),
        (C_NAVY,       "02", "A retention flywheel",         "Every buyer re-buys. CAC compounds down, customer lifetime value up."),
        (C_TERRACOTTA, "03", "ROAS 3.2\u00d7 in 90 days",        "From \u2013\u20b913.5L/mo margin loss to +\u20b91 Cr+ contribution margin."),
    ]
    cy2 = Inches(2.95)
    for col, num, h1, h2 in commits:
        ell(s, Inches(0.6), cy2, Inches(0.52), Inches(0.52), col)
        txt(s, num, Inches(0.6), cy2 + Inches(0.1), Inches(0.52), Inches(0.32),
            size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, ml=0, font=FONT_HDR)
        txt(s, h1, Inches(1.3), cy2, Inches(6.8), Inches(0.32),
            size=13.5, bold=True, color=C_TEXT, ml=0, font=FONT_HDR)
        txt(s, h2, Inches(1.3), cy2 + Inches(0.34), Inches(6.8), Inches(0.28),
            size=9.5, italic=True, color=C_GREY, ml=0, font=FONT_BODY)
        cy2 += Inches(1.05)

    bar(s, Inches(0.6), Inches(6.45), Inches(5.5), Inches(0.04), C_FOREST)
    txt(s, "CHAPTER 2 DRIP  ·  SCALE STRATEGY PRESENTATION  ·  2026",
        Inches(0.6), Inches(6.58), Inches(7.0), Inches(0.3),
        size=8.5, color=C_GREY, font=FONT_BODY)

    # Right: solid Forest Green panel for unlocks
    rbox(s, Inches(8.65), Inches(0.4), Inches(4.25), Inches(6.6), C_FOREST)
    box(s, Inches(8.65), Inches(0.4), Inches(4.25), Inches(0.06), C_GOLD)
    txt(s, "WHAT THIS UNLOCKS FOR C2", Inches(8.85), Inches(0.55), Inches(3.85), Inches(0.3),
        size=9, bold=True, color=C_GOLD, ml=0, font=FONT_HDR)

    unlocks = [
        (C_GOLD,       "₹90L \u2192 ₹1.35 Cr",    "Ad spend scaled 1.5\u00d7 with stability"),
        (C_WHITE,      "1.8\u00d7 \u2192 3.2\u00d7",          "Meta & Google blended ROAS target"),
        (C_GOLD,       "₹1.62 Cr \u2192 ₹4.32 Cr", "Gross monthly revenue growth"),
        (C_WHITE,      "\u20138.3% \u2192 +25.0%",         "Contribution Margin 1 (CM1) flip"),
        (C_GOLD,       "<10% \u2192 25.0%+",         "Customer repeat purchase rate"),
        (C_WHITE,      "₹2,200 \u2192 ₹3,200",     "Average Order Value (AOV) via bundle packs"),
    ]
    uy = Inches(1.0)
    for col, val, lbl in unlocks:
        # nested panel cards
        box(s, Inches(8.65), uy, Inches(4.25), Inches(0.88), C_FOREST, C_BORDER, 0.4)
        txt(s, val, Inches(8.85), uy + Inches(0.08),
            Inches(3.85), Inches(0.38), size=14.5, bold=True, color=col, ml=0, font=FONT_HDR)
        txt(s, lbl, Inches(8.85), uy + Inches(0.48),
            Inches(3.85), Inches(0.28), size=8.5, italic=True, color=C_GOLD, ml=0, font=FONT_BODY)
        uy += Inches(0.92)

    footer(s, 8)


# ── BUILD ─────────────────────────────────────────────────────────────
def build():
    print("Building PREMIUM BRIGHT pitch deck …")
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]
    s01(prs, blank); s02(prs, blank); s03(prs, blank); s04(prs, blank)
    s05(prs, blank); s06(prs, blank); s07(prs, blank); s08(prs, blank)
    out = os.path.join(os.getcwd(), "Chapter2_Drip_Scale_Strategy.pptx")
    prs.save(out)
    print(f"✅  {len(prs.slides)} slides → {out}")

if __name__ == "__main__":
    build()

